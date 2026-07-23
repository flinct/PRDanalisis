from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

ROOT = Path(r"C:/Users/MyBook SAGA 12/Desktop/PRDanalisis/presentation")
DATE = "2026-07-22"

BG = RGBColor(11, 16, 32)
PANEL = RGBColor(18, 25, 55)
TEXT = RGBColor(238, 242, 255)
MUTED = RGBColor(168, 176, 211)
ACCENT = RGBColor(245, 158, 11)
LINE = RGBColor(44, 57, 111)

MONTHLY = {
    "complete": {"version": "v2.7.0", "tickets": 44, "summary": ["42 closed items", "2 items still in testing", "Release mostly complete"]},
    "ongoing": {"version": "v2.8.0", "tickets": 73, "summary": ["10 in progress", "6 risk items (test failed/on hold)", "Bug-heavy delivery scope"]},
    "next": {"version": "v2.9.0", "tickets": 39, "summary": ["38 items still planning/spec", "1 item already active", "Need shortlist before build wave"]},
}

WEEKLY = {
    "current_week": [
        "Close remaining v2.7.0 testing items",
        "Triage v2.8.0 failed/on-hold items",
        "Push active v2.8.0 build/testing flow",
    ],
    "last_week": [
        "v2.7.0 reached mostly closed state",
        "v2.8.0 active scope stayed bug-heavy",
        "v2.9.0 backlog kept in planning stage",
    ],
}

PEOPLE = {
    "Tech": ["Naftal", "Agung", "Dany", "Aprizal"],
    "Product": ["Aryo", "Eva", "Atik"],
}

PERSON_NOTES = {
    "Naftal": ["Lead engineering execution across active versions", "Focus: delivery blockers, code path decisions, release readiness"],
    "Agung": ["Support engineering build and bug-fix stream", "Focus: active tickets and implementation follow-up"],
    "Dany": ["Track product-engineering alignment for current scope", "Focus: requirement clarity, milestone sync, progress visibility"],
    "Aprizal": ["Support engineering task throughput and fixes", "Focus: assigned development items and handoff readiness"],
    "Aryo": ["Drive product priorities and scope trade-offs", "Focus: version sequencing and stakeholder alignment"],
    "Eva": ["Support product documentation and follow-up", "Focus: open requirement details and progress updates"],
    "Atik": ["Support QA/product coordination", "Focus: testing feedback, blockers, and next-step clarity"],
}


def bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def txtbox(slide, x, y, w, h, text, size=24, bold=False, color=TEXT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def bullets(slide, x, y, w, h, items, size=20, color=TEXT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def panel(slide, x, y, w, h):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = LINE
    return shape


def card(slide, x, y, w, h, title, value, note=""):
    panel(slide, x, y, w, h)
    txtbox(slide, x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), Inches(0.3), title, 12, False, MUTED)
    txtbox(slide, x + Inches(0.2), y + Inches(0.45), w - Inches(0.4), Inches(0.5), value, 24, True, TEXT)
    if note:
        txtbox(slide, x + Inches(0.2), y + Inches(1.0), w - Inches(0.4), Inches(0.35), note, 10, False, MUTED)


def add_table(slide, x, y, w, h, headers, rows):
    table = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w, h).table
    for i, head in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = head
        cell.fill.solid()
        cell.fill.fore_color.rgb = PANEL
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = TEXT
    for r_i, row in enumerate(rows, start=1):
        for c_i, val in enumerate(row):
            cell = table.cell(r_i, c_i)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT
    return table


def new_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_title_slide(prs, title, left_cards):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    txtbox(s, Inches(0.6), Inches(0.5), Inches(9), Inches(0.8), title, 26, True)
    txtbox(s, Inches(0.6), Inches(1.2), Inches(8), Inches(0.4), f'Meeting date: {DATE} | Generated date: {DATE}', 14, False, MUTED)
    for i, (head, value, note) in enumerate(left_cards):
        card(s, Inches(0.6 + 4.1 * i), Inches(2.0), Inches(3.8), Inches(1.6), head, value, note)


def add_person_slide(prs, deck_type, team, name):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    txtbox(s, Inches(0.7), Inches(0.45), Inches(8), Inches(0.6), f'{name} — {team}', 24, True)
    txtbox(s, Inches(0.7), Inches(1.0), Inches(6), Inches(0.35), f'{deck_type} owner slide', 12, False, MUTED)
    panel(s, Inches(0.7), Inches(1.5), Inches(12.0), Inches(4.8))
    txtbox(s, Inches(1.0), Inches(1.8), Inches(3), Inches(0.4), 'Current focus', 16, True, ACCENT)
    bullets(s, Inches(1.0), Inches(2.3), Inches(11), Inches(2.8), PERSON_NOTES[name], 18)
    txtbox(s, Inches(1.0), Inches(5.5), Inches(3), Inches(0.35), 'Version context', 16, True, ACCENT)
    if deck_type == 'Monthly':
        items = [
            f"Complete version: {MONTHLY['complete']['version']} ({MONTHLY['complete']['tickets']} tickets)",
            f"Ongoing version: {MONTHLY['ongoing']['version']} ({MONTHLY['ongoing']['tickets']} tickets)",
            f"Next planned version: {MONTHLY['next']['version']} ({MONTHLY['next']['tickets']} tickets)",
        ]
    else:
        items = [
            'This week: ' + WEEKLY['current_week'][0],
            'This week: ' + WEEKLY['current_week'][1],
            'Last week: ' + WEEKLY['last_week'][0],
        ]
    bullets(s, Inches(1.0), Inches(5.9), Inches(11), Inches(0.9), items, 16, MUTED)


def monthly():
    prs = new_prs()
    add_title_slide(prs, 'SatuInbox Product Progress — Monthly', [
        ('Complete Version', MONTHLY['complete']['version'], f"{MONTHLY['complete']['tickets']} tickets"),
        ('Ongoing Version', MONTHLY['ongoing']['version'], f"{MONTHLY['ongoing']['tickets']} tickets"),
        ('Next Planned Version', MONTHLY['next']['version'], f"{MONTHLY['next']['tickets']} tickets"),
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    txtbox(s, Inches(0.6), Inches(0.4), Inches(7), Inches(0.6), 'Monthly Summary', 24, True)
    card(s, Inches(0.6), Inches(1.2), Inches(3.8), Inches(1.5), 'Complete Ver Summary', MONTHLY['complete']['version'], f"{MONTHLY['complete']['tickets']} total tickets")
    card(s, Inches(4.75), Inches(1.2), Inches(3.8), Inches(1.5), 'Ongoing Ver Summary', MONTHLY['ongoing']['version'], f"{MONTHLY['ongoing']['tickets']} total tickets")
    card(s, Inches(8.9), Inches(1.2), Inches(3.8), Inches(1.5), 'Next Planned Ver Total Ticket', str(MONTHLY['next']['tickets']), MONTHLY['next']['version'])
    bullets(s, Inches(0.9), Inches(3.2), Inches(3.6), Inches(2.3), MONTHLY['complete']['summary'], 18)
    bullets(s, Inches(5.0), Inches(3.2), Inches(3.6), Inches(2.3), MONTHLY['ongoing']['summary'], 18)
    bullets(s, Inches(9.15), Inches(3.2), Inches(3.2), Inches(2.3), MONTHLY['next']['summary'], 18)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    txtbox(s, Inches(0.6), Inches(0.4), Inches(5), Inches(0.6), 'Version Timeline', 24, True)
    add_table(s, Inches(0.6), Inches(1.2), Inches(12.0), Inches(2.4), ['Version', 'Status', 'Deadline', 'Actual'], [
        ['v2.7.0', 'Released / Mostly Closed', 'Belum tersedia', 'Belum tersedia di source fetch'],
        ['v2.8.0', 'In Progress', 'Estimated', 'Estimated'],
        ['v2.9.0', 'Specification / Planning', 'Estimated', 'Estimated'],
    ])
    txtbox(s, Inches(0.7), Inches(4.2), Inches(11.5), Inches(0.5), 'Rule: source has no date field, so unreleased versions stay Estimated.', 14, False, MUTED)

    for team, names in PEOPLE.items():
        for name in names:
            add_person_slide(prs, 'Monthly', team, name)

    out = ROOT / 'monthly' / f'{DATE}_satuinbox_monthly_progress.pptx'
    prs.save(out)
    return out, len(prs.slides)


def weekly():
    prs = new_prs()
    add_title_slide(prs, 'SatuInbox Product Progress — Weekly', [
        ('Closed Version', MONTHLY['complete']['version'], ''),
        ('Ongoing Version', MONTHLY['ongoing']['version'], ''),
        ('Next Planned Version', MONTHLY['next']['version'], ''),
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    txtbox(s, Inches(0.6), Inches(0.4), Inches(7), Inches(0.6), 'Weekly Progress Summary', 24, True)
    card(s, Inches(0.6), Inches(1.2), Inches(5.8), Inches(1.5), 'Minggu Berjalan', 'Current Week Progress', 'Focus minggu ini')
    card(s, Inches(6.8), Inches(1.2), Inches(5.8), Inches(1.5), 'Minggu Lalu', 'Last Week Progress', 'Hasil minggu lalu')
    bullets(s, Inches(0.9), Inches(3.0), Inches(5.2), Inches(2.4), WEEKLY['current_week'], 18)
    bullets(s, Inches(7.1), Inches(3.0), Inches(5.0), Inches(2.4), WEEKLY['last_week'], 18)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    txtbox(s, Inches(0.6), Inches(0.4), Inches(5), Inches(0.6), 'Problem / Blocking', 24, True)
    bullets(s, Inches(0.8), Inches(1.5), Inches(11.4), Inches(3.2), [
        'v2.8.0 test failed: 4 items',
        'v2.8.0 on hold: 2 items',
        'v2.7.0 in testing: 2 items',
    ], 22)

    for team, names in PEOPLE.items():
        for name in names:
            add_person_slide(prs, 'Weekly', team, name)

    out = ROOT / 'weekly' / f'{DATE}_satuinbox_weekly_progress.pptx'
    prs.save(out)
    return out, len(prs.slides)


def _self_check():
    monthly_out, monthly_slides = monthly()
    weekly_out, weekly_slides = weekly()
    assert monthly_out.exists() and monthly_out.stat().st_size > 0
    assert weekly_out.exists() and weekly_out.stat().st_size > 0
    assert monthly_slides == 10, monthly_slides
    assert weekly_slides == 10, weekly_slides
    return monthly_out, weekly_out, monthly_slides, weekly_slides


if __name__ == '__main__':
    monthly_out, weekly_out, monthly_slides, weekly_slides = _self_check()
    print(monthly_out)
    print(weekly_out)
    print(f'monthly_slides={monthly_slides}')
    print(f'weekly_slides={weekly_slides}')
